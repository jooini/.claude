#!/usr/bin/env python3
"""Render the presentation pipeline from registry declaration to concrete outputs."""

from __future__ import annotations

import argparse
import json
import re
import shutil
import subprocess
import tempfile
from pathlib import Path


ROOT = Path(__file__).resolve().parents[1]
PIPELINE_PATH = ROOT / "registry" / "presentation-pipeline.json"
DEFAULT_PRESENTATION_DIR = ROOT / "cache" / "presentations"


def load_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def run(cmd: list[str]) -> tuple[int, str, str]:
    process = subprocess.run(cmd, cwd=ROOT, capture_output=True, text=True)
    return process.returncode, process.stdout, process.stderr


def find_tool(name: str) -> bool:
    return shutil.which(name) is not None


def parse_mermaid_blocks(markdown_text: str) -> list[str]:
    pattern = re.compile(r"^```\s*mermaid\s*\n(.*?)^```", re.MULTILINE | re.DOTALL)
    return [block.strip() for block in pattern.findall(markdown_text) if block.strip()]


def resolve_output_path(pipeline_id: str, consumer: dict) -> Path:
    input_path = ROOT / str(consumer["input"])
    output = consumer.get("output")
    if output:
        return ROOT / output

    if consumer.get("type") == "presentation-deck" and consumer.get("format") == "pptx":
        return DEFAULT_PRESENTATION_DIR / f"{pipeline_id}.pptx"

    return input_path


def render_with_marp(input_path: Path, output_path: Path) -> tuple[bool, str | None]:
    if not find_tool("marp"):
        return False, "marp command not found"

    output_path.parent.mkdir(parents=True, exist_ok=True)
    code, stdout, stderr = run(
        [
            "marp",
            "--allow-local-files",
            "--pptx",
            "--output",
            str(output_path),
            str(input_path),
        ]
    )
    if code != 0:
        error = (stderr or stdout or "marp failed without stderr/stdout").strip()
        return False, error[:500]

    return True, None


def render_with_mmdc_python_pptx(input_path: Path, output_path: Path) -> tuple[bool, str | None]:
    if not find_tool("mmdc"):
        return False, "mermaid-cli (mmdc) command not found"

    try:
        from pptx import Presentation
        from pptx.util import Inches
    except Exception as error:
        return False, f"python-pptx import failed: {error}"

    markdown_text = input_path.read_text(encoding="utf-8")
    blocks = parse_mermaid_blocks(markdown_text)
    if not blocks:
        return False, "no mermaid blocks found for fallback conversion"

    output_path.parent.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory() as tempdir:
        tempdir_path = Path(tempdir)
        image_paths: list[Path] = []

        for index, mermaid_block in enumerate(blocks, start=1):
            mermaid_input = tempdir_path / f"diagram_{index:02d}.mmd"
            image_output = tempdir_path / f"diagram_{index:02d}.png"
            mermaid_input.write_text(mermaid_block, encoding="utf-8")

            code, stdout, stderr = run(
                [
                    "mmdc",
                    "-i",
                    str(mermaid_input),
                    "-o",
                    str(image_output),
                    "-w",
                    "1366",
                    "-H",
                    "768",
                    "-b",
                    "white",
                ]
            )
            if code != 0:
                message = (stderr or stdout or "mmdc failed without details").strip()
                return False, f"mmdc failed for diagram #{index}: {message[:500]}"

            image_paths.append(image_output)

        presentation = Presentation()
        presentation.slide_width = Inches(13.333)
        presentation.slide_height = Inches(7.5)
        blank = presentation.slide_layouts[6]

        for image in image_paths:
            slide = presentation.slides.add_slide(blank)
            slide.shapes.add_picture(
                str(image),
                Inches(0),
                Inches(0),
                width=presentation.slide_width,
                height=presentation.slide_height,
            )

        presentation.save(str(output_path))

    return True, None


def has_presentation_runner() -> tuple[bool, str | None]:
    if find_tool("marp"):
        return True, None
    if find_tool("mmdc"):
        try:
            from pptx import Presentation  # noqa: F401
            from pptx.util import Inches  # noqa: F401

            return True, None
        except Exception:
            return False, "fallback requires python-pptx with mmdc"
    return False, "no marp or mmdc+python-pptx runner available"


def run_consumer(pipeline_id: str, consumer: dict) -> tuple[bool, str]:
    consumer_type = consumer.get("type")
    input_path = ROOT / str(consumer["input"])
    if not input_path.exists():
        return False, f"input missing: {consumer.get('input')}"

    output_path = resolve_output_path(pipeline_id, consumer)

    if consumer_type == "markdown-document":
        if output_path == input_path:
            return True, f"markdown up-to-date: {output_path}"
        output_path.parent.mkdir(parents=True, exist_ok=True)
        shutil.copy2(input_path, output_path)
        return True, f"markdown copied: {output_path}"

    if consumer_type == "presentation-deck" and consumer.get("format") == "pptx":
        ok, error = render_with_marp(input_path, output_path)
        if ok:
            return True, f"pptx generated via marp: {output_path}"

        ok, fallback_error = render_with_mmdc_python_pptx(input_path, output_path)
        if ok:
            return True, f"pptx generated via mermaid-cli+python-pptx: {output_path}"

        return False, fallback_error or error or "no pptx renderer available"

    return False, f"unsupported consumer: {consumer_type}"


def run_pipeline(pipeline: dict) -> tuple[int, list[str], list[str]]:
    pipeline_id = pipeline["pipeline_id"]
    ok_messages: list[str] = []
    error_messages: list[str] = []

    for consumer in pipeline.get("consumers", []):
        success, message = run_consumer(pipeline_id, consumer)
        if success:
            ok_messages.append(message)
        else:
            error_messages.append(message)

    return len(error_messages), ok_messages, error_messages


def check_pipeline(pipeline: dict) -> tuple[int, list[str], list[str]]:
    pipeline_id = pipeline["pipeline_id"]
    errors: list[str] = []
    notes: list[str] = [
        f"pipeline={pipeline_id}",
        f"consumers={len(pipeline.get('consumers', []))}",
    ]

    for consumer in pipeline.get("consumers", []):
        output_path = resolve_output_path(pipeline_id, consumer)
        if consumer.get("type") == "markdown-document":
            notes.append(f"markdown-output={output_path.relative_to(ROOT)}")
            continue
        if consumer.get("type") == "presentation-deck" and consumer.get("format") == "pptx":
            available, reason = has_presentation_runner()
            notes.append(f"pptx-runner-available={available}")
            if not available:
                errors.append(reason or "pptx runner unavailable")
            continue

        errors.append(f"unsupported consumer: {consumer.get('type')}")

    return len(errors), notes, errors


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser()
    parser.add_argument("--pipeline", default=str(PIPELINE_PATH), help="pipeline json path")
    parser.add_argument("--check", action="store_true", help="validate availability only")
    parser.add_argument("--list", action="store_true", help="list resolved outputs")
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    pipeline_path = Path(args.pipeline)
    pipeline = load_json(pipeline_path)

    if args.list:
        for consumer in pipeline.get("consumers", []):
            consumer_id = consumer.get("id", "unknown")
            output = resolve_output_path(pipeline["pipeline_id"], consumer)
            print(f"- {consumer_id}: {consumer.get('type')} -> {output.relative_to(ROOT)}")
        return 0

    if args.check:
        error_count, notes, errors = check_pipeline(pipeline)
        print("OK" if not errors else "FAIL")
        for note in notes:
            print(note)
        for message in errors:
            print(f"ERROR: {message}")
        return 1 if error_count else 0

    error_count, ok_messages, error_messages = run_pipeline(pipeline)
    for message in ok_messages:
        print(f"OK: {message}")
    for message in error_messages:
        print(f"ERROR: {message}")

    return 1 if error_count else 0


if __name__ == "__main__":
    raise SystemExit(main())
